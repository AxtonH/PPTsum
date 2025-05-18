# ppt_summary_app_with_chatbot.py
"""
Streamlit PPT Summarizer, Design Assistant & Conversational Chatbot
───────────────────────────────────────────────────────────────────
• Upload a PowerPoint → get deck summary, context, per‑slide insights
• Suggests layouts, charts, visuals, colour schemes **(kept from v1)**
• NEW: built‑in chatbot that answers user questions about the deck
"""

# ── std‑lib ─────────────────────────────────────────────────────────
from io import BytesIO
from typing import List, Dict, Any, Tuple
import os, json, time, requests
# ── 3rd‑party ──────────────────────────────────────────────────────
import streamlit as st
from pptx import Presentation
from dotenv import load_dotenv
from openai import AzureOpenAI
import hydralit_components as hc
# ──────────────────────────────────────────────────────────────────

# ╭─────────────────────  Azure OpenAI config  ─────────────────────╮
load_dotenv()
AZURE_API_KEY  = os.getenv("AZURE_OPENAI_API_KEY")
AZURE_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT")
API_VERSION    = os.getenv("AZURE_OPENAI_API_VERSION", "2024-02-15-preview")
CHAT_MODEL     = os.getenv("AZURE_OPENAI_DEPLOYMENT_NAME", "gpt-4o-mini")
IMAGE_MODEL    = os.getenv("AZURE_OPENAI_IMAGE_DEPLOYMENT_NAME", "dall-e-3")
DEBUG_MODE     = os.getenv("DEBUG_MODE", "False").lower() == "true"
# ╰─────────────────────────────────────────────────────────────────╯

# ──────────────────── Azure helper functions ─────────────────────
@st.cache_resource(show_spinner=False)
def get_client():
    if not AZURE_API_KEY or not AZURE_ENDPOINT:
        st.error("🚨 Missing Azure credentials. Set env vars.")
        st.stop()
    return AzureOpenAI(
        api_key        = AZURE_API_KEY,
        azure_endpoint = AZURE_ENDPOINT,
        api_version    = API_VERSION,
    )

@st.cache_data(show_spinner=False)
def smoke_test():
    url = f"{AZURE_ENDPOINT.rstrip('/')}/openai/deployments?api-version={API_VERSION}"
    resp = requests.get(url, headers={"api-key": AZURE_API_KEY})
    if resp.status_code != 200:
        st.error(f"🚨 Endpoint check failed [{resp.status_code}]: {resp.text}")
        st.stop()
    return resp.json().get("data", [])

# unified chat helper with tiny retry
def chat(system: str, user: str, temperature: float = 0.3, max_attempts: int = 2) -> str:
    client = get_client()
    for attempt in range(max_attempts):
        try:
            resp = client.chat.completions.create(
                model       = CHAT_MODEL,
                messages    = [
                    {"role": "system", "content": system},
                    {"role": "user",   "content": user},
                ],
                temperature = temperature,
            )
            return resp.choices[0].message.content.strip()
        except Exception:
            if attempt + 1 == max_attempts:
                raise
            time.sleep(1.5)

# ────────────────────── Slide extraction ––––––––––––––––––––––––––
def extract_slide_data(ppt_io: BytesIO) -> List[Dict[str, Any]]:
    prs = Presentation(ppt_io)
    data = []
    for i, slide in enumerate(prs.slides, 1):
        text = [sh.text for sh in slide.shapes if hasattr(sh, "text") and sh.text.strip()]
        notes = ""
        if getattr(slide, "notes_slide", None):
            for sh in slide.notes_slide.shapes:
                if hasattr(sh, "text") and sh.text.strip():
                    notes += sh.text + "\n"
        data.append({
            "slide_number": i,
            "slide_text": "\n".join(text),
            "notes_text": notes.strip(),
            "has_notes": bool(notes.strip()),
        })
    return data

# ────────────────────── Visual‑prompt helpers (unchanged) ─────────
def generate_enhanced_image_prompt(slide_summary: str, context: Dict[str, str]) -> str:
    system = (
        "Act like a design consultant. Suggest creative, non‑generic visuals "
        "for high‑level/government audiences. Return only the prompt text."
    )
    user = (
        f"Create an image prompt for this slide:\n\n{slide_summary}\n\n"
        f"Context → Topic: {context.get('topic','–')}; Region: {context.get('region','–')}; "
        f"Purpose: {context.get('purpose','–')}. Make it region‑specific."
    )
    return chat(system, user, 0.7)

def generate_second_image_option(slide_summary: str, context: Dict[str, str], first_prompt: str) -> str:
    system = (
        "Act like a design consultant. Suggest an ALTERNATIVE visual distinct from the first one "
        "for the same high‑level audience. Return only the prompt."
    )
    user = (
        f"Alt image prompt for slide:\n\n{slide_summary}\n\nContext → Topic: {context.get('topic','–')}; "
        f"Region: {context.get('region','–')}; Purpose: {context.get('purpose','–')}."
        f" Previous suggestion: {first_prompt}."
    )
    return chat(system, user, 0.8)

def suggest_chart_type(slide_content: str) -> str:
    system = (
        "You are a data‑viz expert. If a chart suits this content, name the type and why. "
        "Else say no chart needed and suggest another visual."
    )
    return chat(system, slide_content, 0.4)

def suggest_layout(slide_content: str) -> str:
    system = (
        "Consulting‑style slide coach: recommend optimal layout/structure for clarity and impact."
    )
    return chat(system, slide_content, 0.4)

# ────────────────────── Context & summaries ───────────────────────
@st.cache_data(show_spinner=False, ttl=3600)
def identify_context(slides: List[Dict[str, Any]]) -> Dict[str, str]:
    block = []
    for d in slides[:min(5, len(slides))]:
        b = f"Slide {d['slide_number']} content:\n{d['slide_text']}"
        if d["has_notes"]:
            b += f"\n\nNotes:\n{d['notes_text']}"
        block.append(b)
    raw = chat(
        "Return JSON with keys topic, region, purpose inferred from the slides below.",
        "\n\n".join(block),
    )
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        return {k: "Unknown" for k in ("topic", "region", "purpose")}

@st.cache_data(show_spinner=False)
def summarize_deck(slides: List[Dict[str, Any]]) -> Tuple[str, Dict[str, str], List[str], List[str], List[str], List[str]]:
    ctx = identify_context(slides)
    full = []
    for d in slides:
        segment = f"Slide {d['slide_number']} content:\n{d['slide_text']}"
        if d["has_notes"]:
            segment += f"\n\nNotes:\n{d['notes_text']}"
        full.append(segment)
    deck_summary = chat("Expert analyst. Summarise deck <100 words.", "\n\n".join(full))
    slide_sums, design_notes, layouts, charts = [], [], [], []
    for d in slides:
        prompt = f"Slide content: {d['slide_text']}" + (f"\nNotes: {d['notes_text']}" if d["has_notes"] else "")
        resp = chat("Give 1.Summary ≤40w 2.Design Tips ≤40w", prompt)
        p = resp.split("2.")
        slide_sums.append(p[0].replace("1.", "").strip() if p else resp)
        design_notes.append(p[1].strip() if len(p) > 1 else "—")
        layouts.append(suggest_layout(prompt))
        charts.append(suggest_chart_type(prompt))
    return deck_summary, ctx, slide_sums, design_notes, layouts, charts

# ───────────────────── Chatbot corpus helper ──────────────────────
def build_chat_corpus(deck_sum: str, ctx: Dict[str, str],
                      slide_sums: List[str], design: List[str],
                      layouts: List[str], charts: List[str],
                      slides: List[Dict[str, Any]]) -> str:
    bits = [f"Presentation context: {json.dumps(ctx)}", f"Deck summary: {deck_sum}"]
    for i, d in enumerate(slides):
        block = [f"Slide {d['slide_number']}", f"Content: {d['slide_text']}"]
        if d["has_notes"]:
            block.append(f"Notes: {d['notes_text']}")
        block.extend([
            f"Summary: {slide_sums[i]}",
            f"Design tips: {design[i]}",
            f"Layout suggestion: {layouts[i]}",
            f"Chart suggestion: {charts[i]}",
        ])
        bits.append("\n".join(block))
    return "\n\n".join(bits)

# ─────────────────────────── Streamlit UI ─────────────────────────
def main():
    st.set_page_config(
        "PPT Assistant", 
        layout="wide",
        page_icon="📊"
    )
    
    # Custom CSS for better styling
    st.markdown("""
    <style>
        .sidebar .sidebar-content {
            background-color: #f8f9fa;
        }
        .status-box {
            border-radius: 0.5rem;
            padding: 1rem;
            margin-bottom: 1rem;
            background-color: #f0f2f6;
        }
        .status-header {
            font-size: 1.1rem;
            font-weight: 600;
            margin-bottom: 0.5rem;
            color: #2c3e50;
        }
        .status-value {
            font-family: monospace;
            word-break: break-word;
        }
        .slide-nav {
            margin-bottom: 2rem;
        }
        .chat-message {
            border-radius: 0.5rem;
            padding: 0.75rem;
            margin: 0.5rem 0;
        }
        .user-message {
            background-color: #e3f2fd;
        }
        .assistant-message {
            background-color: #f5f5f5;
        }
    </style>
    """, unsafe_allow_html=True)

    # Main title with icon
    st.title("📊 PPT Summarizer, Design Assistant & Chatbot")
    st.caption("Upload a PowerPoint deck to get summaries, design suggestions, and chat with an AI assistant")

    # Enhanced sidebar with Azure status
    with st.sidebar:
        st.subheader("🔌 Azure OpenAI Status")
        
        # Connection status indicator
        status_container = st.container()
        
        if AZURE_API_KEY and AZURE_ENDPOINT:
            try:
                deployments = smoke_test()
                status_container.success("✅ Connected to Azure OpenAI")
                
                # Display deployment info in expander
                with st.expander("🔍 Deployment Details"):
                    st.markdown(f"**API Version:** `{API_VERSION}`")
                    st.markdown(f"**Chat Model:** `{CHAT_MODEL}`")
                    if DEBUG_MODE:
                        st.markdown("**Debug Info:**")
                        st.code(f"Endpoint: {AZURE_ENDPOINT}")
                        if deployments:
                            st.markdown("**Available Deployments:**")
                            for dep in deployments:
                                st.code(f"{dep.get('model')} (id: {dep.get('id')})")
            except Exception as e:
                status_container.error(f"❌ Connection failed: {str(e)}")
        else:
            status_container.error("❌ Missing Azure credentials")
        
        # File upload section in sidebar
        st.divider()
        st.subheader("📂 Upload PPTX")
        ppt = st.file_uploader(
            "Choose a PowerPoint file", 
            type=["pptx"],
            label_visibility="collapsed"
        )
        
        # Debug info (only shown in debug mode)
        if DEBUG_MODE:
            st.divider()
            st.subheader("🔧 Debug Info")
            st.markdown(f"**Session State Keys:**")
            st.write(list(st.session_state.keys()))

    # Main content area
    if ppt is None:
        st.info("⬆️ Upload a PowerPoint deck using the sidebar to get started")
        return

    # Clear session state if new file
    if st.session_state.get("_current_file") != ppt.name:
        st.session_state.clear()
        st.session_state._current_file = ppt.name

    # Extract slides
    slides = extract_slide_data(ppt)
    st.info(f"📑 Found {len(slides)} slides • 📝 {sum(s['has_notes'] for s in slides)} with notes")

    # Analyse button
    if st.button("🚀 Analyse deck", type="primary") or "summaries" in st.session_state:
        if "summaries" not in st.session_state:
            with st.spinner("🔍 Analysing deck with Azure (this may take a minute)..."):
                st.session_state.summaries = summarize_deck(slides)
        deck_sum, ctx, slide_sums, design_notes, layouts, charts = st.session_state.summaries

        # build corpus once
        if "corpus" not in st.session_state:
            st.session_state.corpus = build_chat_corpus(deck_sum, ctx, slide_sums, design_notes, layouts, charts, slides)

        # Presentation context + summary
        with st.expander("📌 Deck Overview", expanded=True):
            cols = st.columns(3)
            cols[0].markdown(f"**📌 Topic:** {ctx.get('topic','–')}")
            cols[1].markdown(f"**🌍 Region:** {ctx.get('region','–')}")
            cols[2].markdown(f"**🎯 Purpose:** {ctx.get('purpose','–')}")
            st.markdown("**📋 Summary:**")
            st.write(deck_sum)
        
        st.divider()

        # Slide navigator
        st.markdown("## Slide Analysis")
        menu = [{"id": i, "label": f"Slide {d['slide_number']}", "icon": "bi-easel"} for i, d in enumerate(slides)]
        sel = hc.nav_bar(menu, sticky_mode="pinned", key="nav")
        idx = sel if isinstance(sel, int) else 0
        s = slides[idx]

        st.subheader(f"📄 Slide {s['slide_number']} Analysis")
        left, right = st.columns([3, 2])
        
        with left:
            with st.container(border=True):
                st.markdown("#### 📝 Summary")
                st.write(slide_sums[idx])
            
            with st.container(border=True):
                st.markdown("#### 🎨 Design Tips")
                st.write(design_notes[idx])
            
            if s["has_notes"]:
                with st.expander("🗒️ Presenter Notes"):
                    st.info(s["notes_text"])
        
        with right:
            with st.container(border=True):
                st.markdown("#### 🖼️ Recommended Layout")
                st.write(layouts[idx])
            
            with st.container(border=True):
                st.markdown("#### 📊 Data Visualization")
                st.write(charts[idx])
            
            with st.expander("📜 Raw Slide Content"):
                st.text(s["slide_text"])
        
        st.divider()

        # Visualisation + colour scheme section
        st.markdown("## 🖌️ Visual Enhancement")
        p_state = f"prompt_state_{s['slide_number']}"
        p1_key = f"prompt1_{s['slide_number']}"
        p2_key = f"prompt2_{s['slide_number']}"
        c_state = f"color_state_{s['slide_number']}"
        c_key   = f"color_{s['slide_number']}"
        
        if p_state not in st.session_state:
            st.session_state[p_state] = False
        if c_state not in st.session_state:
            st.session_state[c_state] = False

        vis_btn_col, color_btn_col = st.columns(2)
        if vis_btn_col.button("💡 Generate Visual Prompts", key=f"vis_btn_{s['slide_number']}"):
            st.session_state[p_state] = True
        if color_btn_col.button("🎨 Suggest Color Scheme", key=f"col_btn_{s['slide_number']}"):
            st.session_state[c_state] = True

        # process prompts
        if st.session_state[p_state] and p1_key not in st.session_state:
            with st.spinner("Generating visual ideas..."):
                enriched = (
                    f"Slide content: {s['slide_text']}\nDesign tips: {design_notes[idx]}\n"
                    f"Chart suggestion: {charts[idx]}"
                )
                p1 = generate_enhanced_image_prompt(enriched, ctx)
                p2 = generate_second_image_option(enriched, ctx, p1)
                st.session_state[p1_key] = p1
                st.session_state[p2_key] = p2
        
        if st.session_state[c_state] and c_key not in st.session_state:
            with st.spinner("Generating color palette..."):
                scheme = chat(
                    "Presentation design expert: give 3‑5 hex colours as bullet list for this context.",
                    f"Context JSON: {json.dumps(ctx)}\nSlide content: {s['slide_text']}",
                    0.4,
                )
                st.session_state[c_key] = scheme

        # display outputs
        if p1_key in st.session_state or c_key in st.session_state:
            tabs = st.tabs(["🖼️ Visual Prompts", "🎨 Color Scheme"])
            
            with tabs[0]:
                if p1_key in st.session_state:
                    cols = st.columns(2)
                    with cols[0]:
                        st.markdown("**Option 1**")
                        st.code(st.session_state[p1_key], language="text")
                    with cols[1]:
                        st.markdown("**Option 2**")
                        st.code(st.session_state[p2_key], language="text")
                else:
                    st.info("Press the 'Generate Visual Prompts' button above")
            
            with tabs[1]:
                if c_key in st.session_state:
                    st.markdown("**Suggested color scheme:**")
                    st.markdown(st.session_state[c_key])
                else:
                    st.info("Press the 'Suggest Color Scheme' button above")
        
        st.divider()

        # Chatbot section
        st.markdown("## 💬 Presentation Chatbot")
        st.caption("Ask questions about the presentation content")
        
        if "messages" not in st.session_state:
            st.session_state.messages = []
        
        # Display chat messages
        for m in st.session_state.messages:
            with st.chat_message(m["role"]):
                st.write(m["content"])
        
        # Chat input
        q = st.chat_input("Ask a question about the presentation...")
        if q:
            st.session_state.messages.append({"role": "user", "content": q})
            with st.chat_message("assistant"):
                with st.spinner("Thinking..."):
                    a = chat(
                        system="Answer strictly from the reference presentation below.",
                        user=q + "\n\nReference:\n" + st.session_state.corpus,
                    )
                    st.write(a)
            st.session_state.messages.append({"role": "assistant", "content": a})

if __name__ == "__main__":
    main()
